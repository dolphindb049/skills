import collections
import collections.abc
from pptx import Presentation

OUTPUT = r"D:\work\.github\skills\develop-skills\references\Skill-Development-Guide-v2.pptx"

prs = Presentation()


def add_title(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for idx, text in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0


add_title(
    "The Complete Guide to Building Skills for DolphinDB",
    "开发分享会 | jrzhang",
)

add_bullets("目录", [
    "1. 简介（Introduction）",
    "2. 基础（Fundamentals）",
    "3. 规划与设计（Planning and Design）",
    "4. 测试与迭代（Testing and Iteration）",
    "5. 分发与共享（Distribution and Sharing）",
    "6. 模式与故障排除（Patterns and Troubleshooting）",
    "7. 资源与结语（Resources and References）",
])

add_bullets("1) 简介：先厘清四个概念", [
    "自行车=Tool（或 MCP）; 把手脚踏=接口协议; 骑车上班=Task; 人=Agent",
    "Skill 是“骑车方法论”，不是任务本身、不是工具本身、也不是 Agent 本体",
    "价值：一次教会，反复复用，减少重复提示成本",
])

add_bullets("2) 基础：渐进式加载架构", [
    "Level 1 Metadata：启动必载，仅 name/description（低 token）",
    "Level 2 Instructions：触发时加载 SKILL.md 主体（中 token）",
    "Level 3+ Resources：按需读取 scripts/references（高容量、低上下文占用）",
    "结论：结构化 Skill 能避免上下文爆炸，同时保留领域深度",
])

add_bullets("3) 规划与设计：从用例倒推", [
    "先定义 2-3 个典型用例：触发词、步骤、工具、交付结果",
    "从真实复杂任务开局（如服务器代理部署/跨网调试）",
    "让 Agent 先跑通，再反向沉淀成 Skill 文件夹",
])

add_bullets("4) 测试与迭代：核心验收线", [
    "验收标准：A 跑通后，B 在同级模型和新环境也能跑通",
    "测试维度：触发测试、流程测试、结果测试",
    "节奏：第一轮最慢，后续迭代快速收敛（执行+复盘产出 Skill）",
])

add_bullets("5) 分发与共享：Doc as Skill", [
    "旧范式：函数接口文档（面向人）",
    "新范式：Skill 文档（面向 Agent 可执行）",
    "组织实践：开发、文档、测试围绕 Skill 统一编排",
])

add_bullets("6) 模式与故障排除", [
    "模式1 顺序编排：多步骤可追踪",
    "模式2 多工具协作：跨 MCP/数据源联动",
    "模式3 迭代式改进：先可用后精化",
    "模式4 情境感知：根据输入动态选工具",
    "模式5 领域智能：沉淀专业 know-how",
])

add_bullets("FICC 示例：从 3-4 小时到分钟级", [
    "链路：拉数据 -> 曲线拟合 -> 调用定价 -> 基准核对",
    "首轮痛点：日计数、曲线插值、字段对齐、空值处理",
    "沉淀结果：`ficc_pricing_example.dos` 模板可复用",
    "收益：稳定性、速度、跨人协作效率显著提升",
])

add_bullets("结语：Skill 是过渡层，也是组织复利层", [
    "每个人的 Skill 都不完美，但可通过共享与评审持续进化",
    "短中期：Skill 是企业级 Agent 落地的核心资产",
    "长期：随着模型内化增强，部分基础 Skill 可自然退场",
])

add_bullets("会后资料", [
    "主指南：develop-skills/SKILL.md",
    "中文翻译：references/The-Complete-Guide-to-Building-Skill-for-Claude.zh.md",
    "在线文档：references/online-meeting-brief.md",
    "示例脚本：scripts/ficc_pricing_example.dos",
])

prs.save(OUTPUT)
print("PPT generated:", OUTPUT)
